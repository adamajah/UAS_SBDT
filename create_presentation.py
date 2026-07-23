from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title_text, subtitle_text, info_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout to custom style
    
    # Background color (Clean white/indigo accent)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Inter'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(79, 70, 229) # Indigo
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.name = 'Inter'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(74, 85, 104) # Muted text
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    # Info List
    info_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(8.0), Inches(2.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    for item in info_list:
        p_info = tf_info.add_paragraph()
        p_info.text = item
        p_info.font.name = 'Inter'
        p_info.font.size = Pt(12)
        p_info.font.color.rgb = RGBColor(113, 128, 150)
        p_info.alignment = PP_ALIGN.CENTER
        
    return slide

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Inter'
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(79, 70, 229) # Indigo
    
    # Content Text Box
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(4.5))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    for i, pt in enumerate(bullet_points):
        # Format text to handle bold parts
        p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
        p.space_after = Pt(12)
        
        # Simple parser for **bold** text in bullets
        parts = pt.split('**')
        for j, part in enumerate(parts):
            run = p.add_run()
            run.text = part
            run.font.name = 'Inter'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(31, 41, 55) # Dark gray
            if j % 2 == 1:
                run.font.bold = True
                
        p.level = 0
        
    return slide

def main():
    prs = Presentation()
    # Set Slide Size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # SLIDE 1: Cover
    slide1 = add_title_slide(
        prs,
        "Sistem Kasir POS Terdistribusi",
        "Implementasi Replikasi Logikal & Fragmentasi Horizontal (SBDT)",
        [
            "Mata Kuliah: Sistem Basis Data Terdistribusi (SBDT)",
            "Penyusun: [Nama Anda] / [NIM Anda]",
            "Teknik Informatika / Sistem Informasi"
        ]
    )
    notes1 = slide1.notes_slide
    notes1.notes_text_frame.text = (
        "Selamat pagi/siang Bapak/Ibu Dosen Penguji.\n\n"
        "Pada hari ini saya akan mempresentasikan proyek tugas akhir SBDT saya yang berjudul "
        "'Rancang Bangun Sistem Kasir Franchise POS Terdistribusi'. Sistem ini dibuat untuk "
        "menyimulasikan arsitektur basis data multi-cabang ritel modern menggunakan teknologi "
        "PostgreSQL dan Docker."
    )
    
    # SLIDE 2: Latar Belakang
    slide2 = add_content_slide(
        prs,
        "Latar Belakang Masalah",
        [
            "**Tantangan Bisnis Franchise**: Memiliki banyak cabang fisik (seperti Jakarta dan Bandung) yang beroperasi serentak.",
            "**Ketergantungan Sentralisasi**: Jika koneksi ke server pusat terputus, operasional kasir cabang akan lumpuh total.",
            "**Latensi Jaringan**: Mengambil data katalog atau menyimpan transaksi ke server pusat yang jauh mengakibatkan antrean kasir melambat.",
            "**Solusi SBDT (Distributed Database)**: Membekali setiap cabang dengan database lokal yang otonom agar operasional kasir tetap berjalan lancar secara offline."
        ]
    )
    notes2 = slide2.notes_slide
    notes2.notes_text_frame.text = (
        "Latar belakang proyek ini adalah masalah konektivitas pada ritel multi-cabang. "
        "Jika sistem database dibuat terpusat (sentralisasi), ketika jaringan internet cabang terputus, "
        "kasir tidak bisa melayani pembeli. Dengan Sistem Basis Data Terdistribusi, kasir cabang "
        "memiliki database lokal sendiri sehingga tetap bisa bertransaksi meski dalam kondisi offline."
    )
    
    # SLIDE 3: Referensi Teori
    slide3 = add_content_slide(
        prs,
        "Referensi Teori SBDT (Pilar Utama)",
        [
            "1. **Otonomi Lokal (Local Autonomy)**: Cabang mengelola transaksi lokal secara mandiri tanpa bergantung pada server pusat.",
            "2. **Replikasi Logikal (Logical Replication)**: Penyebaran asinkron tabel produk dari Pusat ke Cabang menggunakan WAL.",
            "3. **Fragmentasi Horizontal (Horizontal Fragmentation)**: Pembagian baris tabel transaksi berdasarkan cabang (JKT di JKT, BDG di BDG).",
            "4. **Eventual Consistency & Fault Tolerance**: Data transaksi dijamin sinkron saat jaringan pulih, tanpa mengganggu kasir cabang.",
            "5. **Teorema CAP**: Menyeimbangkan Ketersediaan (Availability) dan Toleransi Jaringan (Partition Tolerance)."
        ]
    )
    notes3 = slide3.notes_slide
    notes3.notes_text_frame.text = (
        "Ada 5 referensi teori SBDT utama yang diterapkan di sini. Pertama, otonomi lokal untuk kemandirian node. "
        "Kedua, replikasi logikal satu arah dari pusat ke cabang untuk katalog produk. Ketiga, fragmentasi horizontal "
        "untuk tabel transaksi berdasarkan cabang. Keempat, eventual consistency untuk penanganan kegagalan jaringan. "
        "Dan kelima, kesesuaian dengan Teorema CAP, di mana kami memprioritaskan ketersediaan kasir daripada konsistensi real-time."
    )

    # SLIDE 4: Arsitektur
    slide4 = add_content_slide(
        prs,
        "Arsitektur Jaringan & Topologi (Docker)",
        [
            "**Pusat HQ (db_pusat)**: Menampung skema database `pos_pusat` di port host 5431.",
            "**Cabang Jakarta (db_cabang_jkt)**: Menampung skema database `pos_jkt` di port host 5433.",
            "**Cabang Bandung (db_cabang_bdg)**: Menampung skema database `pos_bdg` di port host 5434.",
            "**Sync Agent**: Daemon pemindah data asinkron berbasis Node.js yang memantau & mengirimkan transaksi non-synced ke pusat setiap 1 detik.",
            "**pgAdmin 4**: Web GUI terpusat untuk inspeksi visual data di semua node (Port 5050)."
        ]
    )
    notes4 = slide4.notes_slide
    notes4.notes_text_frame.text = (
        "Untuk topologi sistem, kami menggunakan Docker Compose. Ada 3 kontainer database PostgreSQL yang saling terhubung "
        "dalam satu virtual network Docker. Replikasi katalog berjalan otomatis secara database-level dari Pusat ke JKT dan BDG. "
        "Sementara data transaksi dari cabang ditarik secara asinkron ke Pusat oleh Sync Agent berbasis Node.js "
        "yang berjalan setiap 1 detik."
    )
    
    # SLIDE 5: Desain Database & Clean Code
    slide5 = add_content_slide(
        prs,
        "Skema Database & Struktur Clean Code",
        [
            "**Tabel Produk (Replikasi)**: Struktur identik di semua cabang, di-update terpusat lalu direplikasikan otomatis.",
            "**Tabel Transaksi (Fragmentasi)**: Terfragmentasi secara horizontal. Menggunakan **UUID** sebagai Primary Key untuk mencegah bentrokan identitas saat digabungkan di database Pusat.",
            "**Struktur Clean Code**: Memisahkan file konfigurasi basis data (`db.js`), database helper query lifecycle (`db_helper.js`), dan route handler modular (`routes/kasir.js`, `routes/dashboard.js`)."
        ]
    )
    notes5 = slide5.notes_slide
    notes5.notes_text_frame.text = (
        "Pada desain skema, tabel produk memiliki struktur yang sama di setiap node. Namun tabel transaksi menggunakan "
        "fragmentasi horizontal. Untuk menghindari bentrokan ID saat transaksi cabang digabungkan di Pusat, "
        "kami menggunakan UUID. Dari sisi pemrograman, aplikasi web ditulis dengan Node.js menggunakan arsitektur clean code "
        "modular untuk memisahkan konfigurasi, utilitas database helper, dan route handler."
    )

    # SLIDE 6: Uji Coba & Demo
    slide6 = add_content_slide(
        prs,
        "Pembahasan Uji Coba & Demo Sistem",
        [
            "**Uji Replikasi**: Menambahkan barang di Pusat, barang langsung muncul di halaman menu kasir Jakarta & Bandung secara instan.",
            "**Uji Fragmentasi**: Transaksi checkout di Kasir JKT hanya tersimpan di database JKT dan database Pusat, database BDG tetap kosong.",
            "**Uji Offline (Fault Tolerance)**: Database pusat dimatikan (`docker stop db_pusat`). Kasir cabang JKT tetap bisa bertransaksi dengan status `is_synced = false`.",
            "**Eventual Consistency**: Ketika pusat dinyalakan kembali, transaksi tertunda otomatis disinkronkan ke pusat dalam waktu 1 detik."
        ]
    )
    notes6 = slide6.notes_slide
    notes6.notes_text_frame.text = (
        "Terakhir, uji coba sistem membuktikan fungsionalitas SBDT berjalan 100%. Kami mendemonstrasikan replikasi logikal "
        "dari Pusat ke Cabang, fragmentasi horizontal di level kasir, serta toleransi kesalahan (mode offline) di mana kasir "
        "cabang tetap dapat melayani penjualan saat database Pusat down, dan data akan secara otomatis tersinkronisasi 1 detik "
        "setelah koneksi database Pusat pulih kembali."
    )

    # SLIDE 7: Kesimpulan
    slide7 = add_content_slide(
        prs,
        "Kesimpulan Hasil Proyek",
        [
            "Sistem kasir POS franchise terdistribusi berhasil dijalankan secara lokal dengan 3 database mandiri di Docker.",
            "Masalah konektivitas lambat atau offline teratasi sempurna melalui konsep **otonomi lokal**.",
            "Sinkronisasi data multi-cabang terjamin aman dan konsisten lewat middleware asinkron 1-detik.",
            "Bebas konflik Primary Key berkat alokasi UUID terdistribusi.",
            "**Terima Kasih** - Sesi Tanya Jawab"
        ]
    )
    notes7 = slide7.notes_slide
    notes7.notes_text_frame.text = (
        "Kesimpulannya, proyek SBDT ini berhasil membuktikan bahwa arsitektur terdistribusi dapat menyelesaikan kendala "
        "konektivitas dan latensi pada bisnis retail multi-cabang. Terima kasih atas perhatian Bapak/Ibu Dosen, "
        "sekarang saya siap menerima pertanyaan atau saran di sesi tanya jawab."
    )
    
    output_path = "slides_presentasi.pptx"
    prs.save(output_path)
    print(f"🎉 Sukses! File presentasi PowerPoint disimpan di: {output_path}")

if __name__ == '__main__':
    main()
